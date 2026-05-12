"""Extract the PNG image from slides 9 + 10 of the Fynd Annual Board Update
pptx. Save to assets/organisation/ as the customer-logo grid backdrop for
/organisation/external."""
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent.parent
PATH = ROOT / "docs/homepage-notes-compilation/Fynd Annual Board Update 2025-2.pptx"
OUT = ROOT / "assets/organisation"
OUT.mkdir(parents=True, exist_ok=True)

p = Presentation(str(PATH))

NAMES = {9: "external-customers-india.png", 10: "external-customers-global.png"}

for idx, slide in enumerate(p.slides, start=1):
    if idx not in NAMES:
        continue
    for sh in slide.shapes:
        if hasattr(sh, "image") and sh.image is not None:
            ext = sh.image.ext  # e.g. 'png'
            blob = sh.image.blob
            target = OUT / NAMES[idx]
            target.write_bytes(blob)
            print(f"slide {idx} → {target} ({len(blob):,} bytes · {ext})")
            break
