"""Inspect slides 9 + 10 of the Fynd Annual Board Update pptx — see what's on them
and figure out how to render them for the External Commercialization page."""
from pptx import Presentation

PATH = "docs/homepage-notes-compilation/Fynd Annual Board Update 2025-2.pptx"
p = Presentation(PATH)

print(f"total slides: {len(p.slides)}")
print()

for idx, slide in enumerate(p.slides, start=1):
    if idx not in (9, 10):
        continue
    print(f"=== slide {idx} ===")
    print(f"  layout: {slide.slide_layout.name}")
    print(f"  shapes ({len(slide.shapes)}):")
    for sh in slide.shapes:
        kind = sh.shape_type
        name = sh.name
        info = ""
        if sh.has_text_frame:
            txt = sh.text_frame.text.strip().replace("\n", " | ")
            info = f"  text={txt[:80]!r}"
        elif hasattr(sh, "image") and sh.image is not None:
            info = f"  image={sh.image.content_type} ext={sh.image.ext}"
        print(f"    [{kind}] {name}{info}")
    print()
